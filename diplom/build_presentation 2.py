#!/usr/bin/env python3
"""
Генерация презентации к защите ВКР (python-pptx).

Берёт оформление из шаблона курсовой «курсовая .pptx» в корне репозитория
(те же темы слайдов, колонтитулы, шрифты), подставляет текст ВКР.

Запуск из каталога diplom:
  ./.venv-docx/bin/python build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
# Имя файла с пробелом перед расширением, как в репозитории
TEMPLATE = REPO_ROOT / "курсовая .pptx"
OUT = Path(__file__).resolve().parent / "VKR_prezentatsiya_Inevatkin.pptx"

# Индексы макетов из шаблона КубГУ: 0 — титул, 1 — заголовок + текст
LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1


def _delete_slide(prs: Presentation, index: int) -> None:
    sld_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(sld_id.rId)
    prs.slides._sldIdLst.remove(sld_id)


def _bullets(slide, items: list[str]) -> None:
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(6)


def _load_from_template() -> Presentation:
    if not TEMPLATE.exists():
        raise FileNotFoundError(
            f"Не найден шаблон «курсовая .pptx»: ожидался файл {TEMPLATE}"
        )
    prs = Presentation(str(TEMPLATE))
    while len(prs.slides) > 0:
        _delete_slide(prs, 0)
    return prs


def main() -> None:
    prs = _load_from_template()

    # --- Титул (макет как в курсовой) ---
    s0 = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    s0.shapes.title.text = (
        "ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА\n"
        "Автоматизация обработки данных статистического моделирования\n"
        "на основе символьной регрессии"
    )
    st = s0.placeholders[1]
    st.text = (
        "Исполнитель: Иневаткин Ю. А.\n"
        "Научный руководитель: канд. техн. наук, доц. Миков А. И.\n\n"
        "Кубанский государственный университет\n"
        "Факультет компьютерных технологий и прикладной математики\n"
        "Направление 02.03.03 — 2026"
    )
    for p in st.text_frame.paragraphs:
        if p.font.size is None or p.font.size < Pt(16):
            p.font.size = Pt(18)

    slides_data: list[tuple[str, list[str]]] = [
        (
            "Актуальность",
            [
                "Результат моделирования часто — таблица чисел, а нужна компактная аналитическая зависимость.",
                "Классическая регрессия: фиксированная структура; SR/GP — поиск структуры и констант.",
                "Ручной перебор .xlsx и метрик плохо воспроизводим — нужен единый конвейер и протокол.",
            ],
        ),
        (
            "Цель и задачи",
            [
                "Цель: программный комплекс для пакетной обработки книг .xlsx и символьной регрессии с единым отчётом.",
                "Задачи: обзор SR/GP; требования; проектирование; реализация на C# (.NET 8); вычислительный эксперимент.",
            ],
        ),
        (
            "Объект и предмет",
            [
                "Объект: процессы обработки табличных выборок «факторы — отклик».",
                "Предмет: методы символьной регрессии и конвейер импорт → эволюция → метрики → отчёт.",
            ],
        ),
        (
            "Символьная регрессия и GP",
            [
                "Формулы из грамматики (операции, глубина дерева, штрафы за сложность и NaN).",
                "Генетическое программирование: деревья как особи, минимизируется fitness.",
                "В отчёте: RMSE, R², hold-out; для benchmark — RMSE и относительная ошибка rel.",
            ],
        ),
        (
            "Автоматизация: вход и выход",
            [
                "Вход: каталог с .xlsx (первый лист), опционально gp_settings.json и base_functions.txt.",
                "Все книги подхватываются автоматически; служебные ~$* пропускаются.",
                "Выход: result.txt — формулы, метрики, кандидаты (MD), блок benchmark при порогах в JSON.",
            ],
        ),
        (
            "Архитектура комплекса",
            [
                "1. Ввод: ClosedXML, фильтрация строк; 2 столбца → 1D, ≥3 → MD.",
                "2. Конфигурация: профили GP и benchmarkThresholds в JSON.",
                "3. Ядро: эвристики 1D (в т.ч. log-fit); MD — GP + baseline, выбор по test RMSE.",
                "4. Отчёт: result.txt и сводка PASS/FAIL.",
            ],
        ),
        (
            "Реализация (ключевое)",
            [
                "Суперпозиция: базовые φ из base_functions.txt.",
                "PASS в benchmark при RMSE ≤ порога или rel ≤ 0,005 — пометки (abs+rel), (abs), (rel).",
                "Подпись «Оценка (BIC)» в логе — штрафованная ошибка эволюции, не классический BIC.",
            ],
        ),
        (
            "Эксперимент",
            [
                "11 контрольных .xlsx с известной зависимостью; пороги в gp_settings.json.",
                "Итог: 11/11 PASS; отдельно разобраны exponential (PASS по rel) и noisy_trig (PASS по abs).",
                "Примеры: линейная, полином, 2·eˣ, sin x, sin(x₁)+x₂²; MD — baseline vs GP.",
            ],
        ),
        (
            "Значимость и развитие",
            [
                "Практика: серия таблиц → один архивируемый протокол; учебные и исследовательские задачи.",
                "Ограничения: числовой формат, чувствительность к настройке GP.",
                "Развитие: декомпозиция Program.cs, тесты, CLI, расширенный отчёт.",
            ],
        ),
        (
            "Спасибо за внимание!",
            ["Готов ответить на вопросы комиссии."],
        ),
    ]

    for title, items in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
        slide.shapes.title.text = title
        tf = slide.shapes.title.text_frame
        tf.paragraphs[0].font.size = Pt(28)
        _bullets(slide, items)

    prs.save(OUT)
    print(f"Шаблон: {TEMPLATE}")
    print(f"Сохранено: {OUT}")


if __name__ == "__main__":
    main()
